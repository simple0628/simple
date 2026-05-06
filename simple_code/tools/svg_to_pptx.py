"""SVG → PPTX 原生形状转换器（直接写 DrawingML XML，参考 ppt-master）

python-pptx 仅用于创建空白演示文稿和保存文件。
所有形状通过直接构造 DrawingML XML 插入，与 ppt-master 完全一致。
"""

import re
import math
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

# 核心常量
EMU_PER_PX = 9525
FONT_PX_TO_HUNDREDTHS_PT = 75  # 1px = 0.75pt = 75 hundredths-of-a-point
ANGLE_UNIT = 60000  # DrawingML 角度单位 = 1/60000 度

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

_shape_id_counter = 100


def _next_id():
    global _shape_id_counter
    _shape_id_counter += 1
    return _shape_id_counter


def px(v):
    return round(v * EMU_PER_PX)


def parse_num(s, default=0.0):
    if s is None: return default
    s = str(s).strip().replace("px", "")
    m = re.match(r"[-+]?[\d.]+", s)
    return float(m.group()) if m else default


def parse_color_hex(s):
    """返回 6 位 HEX 字符串（无 #）"""
    if not s or s.lower() in ("none", "transparent"): return None
    s = s.strip().lstrip("#")
    if len(s) == 3: s = s[0]*2 + s[1]*2 + s[2]*2
    return s[:6] if len(s) >= 6 else None


def get_attr(elem, name, default=None):
    val = elem.get(name)
    if val: return val
    style = elem.get("style", "")
    for part in style.split(";"):
        if ":" in part:
            k, v = part.split(":", 1)
            if k.strip() == name: return v.strip()
    return default


def is_cjk(ch):
    cp = ord(ch)
    return (0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF or
            0x3000 <= cp <= 0x303F or 0xFF00 <= cp <= 0xFFEF or
            0x2E80 <= cp <= 0x2EFF or 0xF900 <= cp <= 0xFAFF)


def estimate_text_width(text, font_size, bold=False):
    w = 0.0
    for ch in text:
        if is_cjk(ch): w += font_size
        elif ch == ' ': w += font_size * 0.3
        elif ch in 'mMwWOQ@': w += font_size * 0.75
        elif ch in 'iIlj1!|.,;:': w += font_size * 0.3
        else: w += font_size * 0.55
    if bold: w *= 1.05
    return w * 1.15


def parse_transform(s):
    tx, ty, sx, sy, rot = 0, 0, 1, 1, 0
    if not s: return tx, ty, sx, sy, rot
    for m in re.finditer(r"(\w+)\(([^)]+)\)", s):
        func, args = m.group(1), [float(x) for x in re.findall(r"[-+]?[\d.]+", m.group(2))]
        if func == "translate":
            tx += args[0]; ty += args[1] if len(args) > 1 else 0
        elif func == "scale":
            sx *= args[0]; sy *= args[1] if len(args) > 1 else args[0]
        elif func == "rotate":
            rot += args[0]
    return tx, ty, sx, sy, rot


# ============================================================
#  DrawingML XML 构建器
# ============================================================

def _fill_xml(color_hex, opacity=1.0):
    """生成 solidFill XML"""
    if not color_hex: return "<a:noFill/>"
    alpha = f'<a:alpha val="{int(opacity * 100000)}"/>' if opacity < 1 else ""
    return f'<a:solidFill><a:srgbClr val="{color_hex}">{alpha}</a:srgbClr></a:solidFill>'


def _gradient_fill_xml(stops, angle_deg=0):
    """生成 gradFill XML"""
    gs_list = []
    for offset, color_hex, opacity in stops:
        pos = int(offset * 100000)
        alpha = f'<a:alpha val="{int(opacity * 100000)}"/>' if opacity < 1 else ""
        gs_list.append(f'<a:gs pos="{pos}"><a:srgbClr val="{color_hex}">{alpha}</a:srgbClr></a:gs>')
    ang = int(angle_deg * ANGLE_UNIT) % (360 * ANGLE_UNIT)
    return (f'<a:gradFill><a:gsLst>{"".join(gs_list)}</a:gsLst>'
            f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')


def _line_xml(color_hex=None, width_px=1, opacity=1.0):
    if not color_hex: return "<a:ln><a:noFill/></a:ln>"
    w = px(width_px)
    fill = _fill_xml(color_hex, opacity)
    return f'<a:ln w="{w}">{fill}</a:ln>'


def _shape_xml(shape_id, name, off_x, off_y, cx, cy, geom_xml, fill_xml_str, line_xml_str, rot=0):
    """构建完整的 p:sp XML"""
    rot_attr = f' rot="{int(rot * ANGLE_UNIT)}"' if rot else ""
    return (
        f'<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="{name}"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm{rot_attr}><a:off x="{off_x}" y="{off_y}"/>'
        f'<a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        f'{geom_xml}{fill_xml_str}{line_xml_str}'
        f'</p:spPr></p:sp>'
    )


def _text_shape_xml(shape_id, off_x, off_y, cx, cy, runs_xml, algn="l", rot=0):
    """构建文本框 p:sp XML（直接写 DrawingML，与 ppt-master 一致）"""
    rot_attr = f' rot="{int(rot * ANGLE_UNIT)}"' if rot else ""
    return (
        f'<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:nvSpPr><p:cNvPr id="{shape_id}" name="TextBox {shape_id}"/>'
        f'<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm{rot_attr}><a:off x="{off_x}" y="{off_y}"/>'
        f'<a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:noFill/><a:ln><a:noFill/></a:ln>'
        f'</p:spPr>'
        f'<p:txBody>'
        f'<a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" anchor="t" anchorCtr="0">'
        f'<a:spAutoFit/></a:bodyPr>'
        f'<a:lstStyle/>'
        f'<a:p><a:pPr algn="{algn}"/>{runs_xml}</a:p>'
        f'</p:txBody></p:sp>'
    )


def _run_xml(text, font_size_px, color_hex="333333", bold=False, font="Microsoft YaHei"):
    """构建单个文本运行 XML"""
    sz = round(font_size_px * FONT_PX_TO_HUNDREDTHS_PT)
    b_attr = ' b="1"' if bold else ""
    fill = _fill_xml(color_hex) if color_hex else ""
    # 转义 XML 特殊字符
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    return (
        f'<a:r><a:rPr lang="zh-CN" sz="{sz}"{b_attr} dirty="0">'
        f'{fill}'
        f'<a:latin typeface="{font}"/>'
        f'<a:ea typeface="Microsoft YaHei"/>'
        f'<a:cs typeface="{font}"/>'
        f'</a:rPr><a:t>{text}</a:t></a:r>'
    )


# ============================================================
#  SVG Path → DrawingML custGeom
# ============================================================

def _tokenize_path(d):
    tokens = re.findall(r"[MmLlCcSsQqTtAaHhVvZz]|[-+]?[\d]*\.?\d+(?:[eE][-+]?\d+)?", d)
    result, cmd, nums = [], None, []
    for t in tokens:
        if t.isalpha():
            if cmd: result.append((cmd, nums))
            cmd, nums = t, []
        else:
            nums.append(float(t))
    if cmd: result.append((cmd, nums))
    return result


def _to_absolute(commands):
    abs_cmds, cx, cy = [], 0, 0
    for cmd, nums in commands:
        if cmd == 'M': cx, cy = nums[0], nums[1]; abs_cmds.append(('M', [cx, cy]))
        elif cmd == 'm': cx += nums[0]; cy += nums[1]; abs_cmds.append(('M', [cx, cy]))
        elif cmd == 'L':
            for i in range(0, len(nums), 2): cx, cy = nums[i], nums[i+1]; abs_cmds.append(('L', [cx, cy]))
        elif cmd == 'l':
            for i in range(0, len(nums), 2): cx += nums[i]; cy += nums[i+1]; abs_cmds.append(('L', [cx, cy]))
        elif cmd == 'H': cx = nums[0]; abs_cmds.append(('L', [cx, cy]))
        elif cmd == 'h': cx += nums[0]; abs_cmds.append(('L', [cx, cy]))
        elif cmd == 'V': cy = nums[0]; abs_cmds.append(('L', [cx, cy]))
        elif cmd == 'v': cy += nums[0]; abs_cmds.append(('L', [cx, cy]))
        elif cmd == 'C':
            for i in range(0, len(nums), 6): abs_cmds.append(('C', nums[i:i+6])); cx, cy = nums[i+4], nums[i+5]
        elif cmd == 'c':
            for i in range(0, len(nums), 6):
                pts = [cx+nums[i], cy+nums[i+1], cx+nums[i+2], cy+nums[i+3], cx+nums[i+4], cy+nums[i+5]]
                abs_cmds.append(('C', pts)); cx, cy = pts[4], pts[5]
        elif cmd in ('Z', 'z'): abs_cmds.append(('Z', []))
    return abs_cmds


def _path_geom_xml(d_attr):
    """SVG path d → (custGeom XML, min_x, min_y, w, h)"""
    commands = _to_absolute(_tokenize_path(d_attr))
    if not commands: return None, 0, 0, 0, 0
    points = []
    for cmd, nums in commands:
        if cmd in ('M', 'L'): points.append((nums[0], nums[1]))
        elif cmd == 'C': points.extend([(nums[i], nums[i+1]) for i in range(0, 6, 2)])
    if not points: return None, 0, 0, 0, 0
    min_x, min_y = min(p[0] for p in points), min(p[1] for p in points)
    max_x, max_y = max(p[0] for p in points), max(p[1] for p in points)
    w, h = max(max_x - min_x, 1), max(max_y - min_y, 1)

    children = []
    for cmd, nums in commands:
        if cmd == 'M': children.append(f'<a:moveTo><a:pt x="{px(nums[0]-min_x)}" y="{px(nums[1]-min_y)}"/></a:moveTo>')
        elif cmd == 'L': children.append(f'<a:lnTo><a:pt x="{px(nums[0]-min_x)}" y="{px(nums[1]-min_y)}"/></a:lnTo>')
        elif cmd == 'C':
            pts = "".join(f'<a:pt x="{px(nums[i]-min_x)}" y="{px(nums[i+1]-min_y)}"/>' for i in range(0, 6, 2))
            children.append(f'<a:cubicBezTo>{pts}</a:cubicBezTo>')
        elif cmd == 'Z': children.append('<a:close/>')

    xml = (f'<a:custGeom xmlns:a="{NS_A}"><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
           f'<a:rect l="l" t="t" r="r" b="b"/>'
           f'<a:pathLst><a:path w="{px(w)}" h="{px(h)}">{"".join(children)}</a:path></a:pathLst>'
           f'</a:custGeom>')
    return xml, min_x, min_y, w, h


# ============================================================
#  主转换器
# ============================================================

# Chrome ID 关键词 — 这些元素不参与动画
_CHROME_TOKENS = {'background', 'bg', 'decoration', 'decorations', 'decor',
                  'header', 'footer', 'chrome', 'watermark', 'pagenumber', 'pagenum'}

def _is_chrome_id(svg_id):
    if not svg_id: return False
    lower = svg_id.lower().replace("-", "").replace("_", "")
    return any(token in lower for token in _CHROME_TOKENS)


class SvgConverter:
    def __init__(self):
        self.gradients = {}
        self.anim_targets = []  # [(shape_id, svg_id), ...]

    def convert(self, slide, svg_str):
        """转换 SVG → slide 形状，返回 anim_targets"""
        root = etree.fromstring(svg_str.encode("utf-8"))
        self._collect_defs(root)
        sp_tree = slide._element.find(f".//{qn('p:spTree')}")
        self._walk(sp_tree, root, 0, 0, 1, 1, depth=0)
        return self.anim_targets

    def _collect_defs(self, root):
        for elem in root.iter():
            tag = self._tag(elem)
            if tag == "linearGradient":
                gid = elem.get("id")
                if not gid: continue
                x1 = self._pct(elem.get("x1", "0")); y1 = self._pct(elem.get("y1", "0"))
                x2 = self._pct(elem.get("x2", "100")); y2 = self._pct(elem.get("y2", "0"))
                angle = math.degrees(math.atan2(y2 - y1, x2 - x1)) % 360
                stops = []
                for stop in elem:
                    if self._tag(stop) == "stop":
                        offset = self._pct(stop.get("offset", "0"))
                        color = parse_color_hex(stop.get("stop-color", "#000000"))
                        opacity = parse_num(stop.get("stop-opacity", "1"))
                        if color: stops.append((offset, color, opacity))
                if stops: self.gradients[gid] = {"stops": stops, "angle": angle}

    def _pct(self, s):
        if "%" in s: return parse_num(s.replace("%", "")) / 100
        return parse_num(s)

    def _tag(self, elem):
        t = elem.tag
        return t.split("}")[1] if "}" in t else t

    def _resolve_fill(self, elem):
        """返回 (fill_xml, opacity)"""
        fill_str = get_attr(elem, "fill", "#000000")
        opacity = parse_num(get_attr(elem, "fill-opacity", "1")) * parse_num(get_attr(elem, "opacity", "1"))

        if not fill_str or fill_str.lower() in ("none", "transparent"):
            return "<a:noFill/>", opacity

        if fill_str.startswith("url("):
            gid = fill_str.replace("url(#", "").replace(")", "").strip()
            if gid in self.gradients:
                grad = self.gradients[gid]
                return _gradient_fill_xml(grad["stops"], grad["angle"]), opacity
            return "<a:noFill/>", opacity

        color = parse_color_hex(fill_str)
        if color:
            return _fill_xml(color, opacity), 1.0  # opacity 已在 fill 里
        return "<a:noFill/>", opacity

    def _resolve_stroke(self, elem):
        stroke = get_attr(elem, "stroke")
        if not stroke or stroke.lower() in ("none", "transparent"):
            return "<a:ln><a:noFill/></a:ln>"
        color = parse_color_hex(stroke)
        w = parse_num(get_attr(elem, "stroke-width", "1"))
        opacity = parse_num(get_attr(elem, "stroke-opacity", "1"))
        return _line_xml(color, w, opacity)

    def _walk(self, sp_tree, elem, tx, ty, sx, sy, depth=0):
        for child in elem:
            tag = self._tag(child)
            ctxt = parse_transform(child.get("transform"))
            ntx = tx + ctxt[0] * sx
            nty = ty + ctxt[1] * sy
            nsx, nsy = sx * ctxt[2], sy * ctxt[3]

            try:
                if tag == "rect": self._rect(sp_tree, child, ntx, nty, nsx, nsy)
                elif tag == "circle": self._circle(sp_tree, child, ntx, nty, nsx, nsy)
                elif tag == "ellipse": self._ellipse(sp_tree, child, ntx, nty, nsx, nsy)
                elif tag == "line": self._line(sp_tree, child, ntx, nty, nsx, nsy)
                elif tag == "text": self._text(sp_tree, child, ntx, nty, nsx, nsy)
                elif tag == "path": self._path(sp_tree, child, ntx, nty, nsx, nsy)
                elif tag == "polygon": self._polygon(sp_tree, child, ntx, nty, nsx, nsy)
                elif tag in ("g", "svg"):
                    svg_id = child.get("id")
                    # 记录顶层非 chrome 组为动画目标
                    if depth == 0 and svg_id and not _is_chrome_id(svg_id):
                        # 记录下一个将分配的 shape_id
                        next_shape_id = _shape_id_counter + 1
                        self._walk(sp_tree, child, ntx, nty, nsx, nsy, depth + 1)
                        # 如果有新形状被添加，记录为动画目标
                        if _shape_id_counter >= next_shape_id:
                            self.anim_targets.append((next_shape_id, svg_id))
                    else:
                        self._walk(sp_tree, child, ntx, nty, nsx, nsy, depth + 1)
            except Exception:
                pass

    # --- 形状转换 ---

    def _rect(self, sp_tree, elem, tx, ty, sx, sy):
        x = parse_num(elem.get("x")) * sx + tx
        y = parse_num(elem.get("y")) * sy + ty
        w = parse_num(elem.get("width")) * sx
        h = parse_num(elem.get("height")) * sy
        if w <= 0 or h <= 0: return

        rx = parse_num(elem.get("rx"))
        if rx > 0:
            adj = min(int(rx / min(w, h) * 100000 * 2), 50000)
            geom = f'<a:prstGeom prst="roundRect" xmlns:a="{NS_A}"><a:avLst><a:gd name="adj" fmla="val {adj}"/></a:avLst></a:prstGeom>'
        else:
            geom = f'<a:prstGeom prst="rect" xmlns:a="{NS_A}"><a:avLst/></a:prstGeom>'

        fill, _ = self._resolve_fill(elem)
        stroke = self._resolve_stroke(elem)
        xml = _shape_xml(_next_id(), "Rect", px(x), px(y), px(w), px(h), geom, fill, stroke)
        sp_tree.append(etree.fromstring(xml))

    def _circle(self, sp_tree, elem, tx, ty, sx, sy):
        cx = parse_num(elem.get("cx")) * sx + tx
        cy = parse_num(elem.get("cy")) * sy + ty
        r = parse_num(elem.get("r")) * sx
        geom = f'<a:prstGeom prst="ellipse" xmlns:a="{NS_A}"><a:avLst/></a:prstGeom>'
        fill, _ = self._resolve_fill(elem)
        stroke = self._resolve_stroke(elem)
        xml = _shape_xml(_next_id(), "Circle", px(cx-r), px(cy-r), px(r*2), px(r*2), geom, fill, stroke)
        sp_tree.append(etree.fromstring(xml))

    def _ellipse(self, sp_tree, elem, tx, ty, sx, sy):
        cx = parse_num(elem.get("cx")) * sx + tx
        cy = parse_num(elem.get("cy")) * sy + ty
        rx = parse_num(elem.get("rx")) * sx
        ry = parse_num(elem.get("ry")) * sy
        geom = f'<a:prstGeom prst="ellipse" xmlns:a="{NS_A}"><a:avLst/></a:prstGeom>'
        fill, _ = self._resolve_fill(elem)
        stroke = self._resolve_stroke(elem)
        xml = _shape_xml(_next_id(), "Ellipse", px(cx-rx), px(cy-ry), px(rx*2), px(ry*2), geom, fill, stroke)
        sp_tree.append(etree.fromstring(xml))

    def _line(self, sp_tree, elem, tx, ty, sx, sy):
        x1 = parse_num(elem.get("x1")) * sx + tx
        y1 = parse_num(elem.get("y1")) * sy + ty
        x2 = parse_num(elem.get("x2")) * sx + tx
        y2 = parse_num(elem.get("y2")) * sy + ty
        stroke_w = parse_num(get_attr(elem, "stroke-width", "1"))
        stroke_color = parse_color_hex(get_attr(elem, "stroke", "#000000"))
        opacity = parse_num(get_attr(elem, "stroke-opacity", "1"))
        min_x, min_y = min(x1, x2), min(y1, y2)
        w = max(abs(x2 - x1), stroke_w)
        h = max(abs(y2 - y1), stroke_w)
        geom = f'<a:prstGeom prst="rect" xmlns:a="{NS_A}"><a:avLst/></a:prstGeom>'
        fill = _fill_xml(stroke_color, opacity) if stroke_color else "<a:noFill/>"
        line = "<a:ln><a:noFill/></a:ln>"
        xml = _shape_xml(_next_id(), "Line", px(min_x), px(min_y), px(w), px(h), geom, fill, line)
        sp_tree.append(etree.fromstring(xml))

    def _text(self, sp_tree, elem, tx, ty, sx, sy):
        x = parse_num(elem.get("x")) * sx + tx
        y = parse_num(elem.get("y")) * sy + ty
        font_size = parse_num(get_attr(elem, "font-size", "18"))
        fill = get_attr(elem, "fill", "#333333")
        color_hex = parse_color_hex(fill) or "333333"
        bold = get_attr(elem, "font-weight", "") in ("bold", "700", "600", "800", "900")
        anchor = get_attr(elem, "text-anchor", "start")
        font = get_attr(elem, "font-family", "Microsoft YaHei")
        font = font.strip("'\"").split(",")[0].strip()

        # 收集文本运行
        runs = []
        if elem.text and elem.text.strip():
            runs.append({"text": elem.text.strip(), "size": font_size, "color": color_hex, "bold": bold, "font": font})
        for child in elem:
            if self._tag(child) == "tspan" and child.text:
                cs = parse_num(get_attr(child, "font-size", str(font_size)))
                cc = parse_color_hex(get_attr(child, "fill", fill)) or color_hex
                cb = get_attr(child, "font-weight", "bold" if bold else "") in ("bold", "700", "600")
                cf = get_attr(child, "font-family", font).strip("'\"").split(",")[0].strip()
                if child.get("x") is not None: x = parse_num(child.get("x")) * sx + tx
                if child.get("y") is not None: y = parse_num(child.get("y")) * sy + ty
                runs.append({"text": child.text.strip(), "size": cs, "color": cc, "bold": cb, "font": cf})
            if child.tail and child.tail.strip():
                runs.append({"text": child.tail.strip(), "size": font_size, "color": color_hex, "bold": bold, "font": font})

        if not runs: return

        full_text = "".join(r["text"] for r in runs)
        text_width = estimate_text_width(full_text, font_size, bold)
        padding = font_size * 0.1
        text_height = font_size * 1.5

        # 定位（与 ppt-master 一致）
        if anchor == "middle":
            box_x = x - text_width / 2 - padding
            algn = "ctr"
        elif anchor == "end":
            box_x = x - text_width - padding
            algn = "r"
        else:
            box_x = x - padding
            algn = "l"

        box_y = y - font_size * 0.85  # ppt-master 的基线偏移
        box_w = text_width + padding * 2
        box_h = text_height + padding

        # 构建 runs XML
        runs_xml = ""
        for r in runs:
            runs_xml += _run_xml(r["text"], r["size"], r["color"], r["bold"], r["font"])

        xml = _text_shape_xml(_next_id(), px(max(0, box_x)), px(max(0, box_y)),
                              px(max(box_w, font_size)), px(box_h), runs_xml, algn)
        sp_tree.append(etree.fromstring(xml))

    def _path(self, sp_tree, elem, tx, ty, sx, sy):
        d = elem.get("d")
        if not d: return
        geom_xml, min_x, min_y, w, h = _path_geom_xml(d)
        if not geom_xml: return
        fill, _ = self._resolve_fill(elem)
        stroke = self._resolve_stroke(elem)
        xml = _shape_xml(_next_id(), "Path", px(min_x * sx + tx), px(min_y * sy + ty),
                         px(w * sx), px(h * sy), geom_xml, fill, stroke)
        sp_tree.append(etree.fromstring(xml))

    def _polygon(self, sp_tree, elem, tx, ty, sx, sy):
        points_str = elem.get("points", "")
        nums = [float(x) for x in re.findall(r"[-+]?[\d.]+", points_str)]
        if len(nums) < 4: return
        pairs = [(nums[i], nums[i+1]) for i in range(0, len(nums) - 1, 2)]
        d = f"M {pairs[0][0]},{pairs[0][1]} " + " ".join(f"L {p[0]},{p[1]}" for p in pairs[1:]) + " Z"
        fake = etree.Element("path")
        fake.set("d", d)
        for attr in ("fill", "stroke", "stroke-width", "fill-opacity", "opacity", "style"):
            v = elem.get(attr)
            if v: fake.set(attr, v)
        self._path(sp_tree, fake, tx, ty, sx, sy)


# ============================================================
#  动画 & 转场 XML 生成（参考 ppt-master）
# ============================================================

TRANSITIONS = {
    "fade":    {"name": "淡入淡出", "duration_ms": 700, "element": "fade",    "attrs": {}},
    "push":    {"name": "推入",     "duration_ms": 500, "element": "push",    "attrs": {"dir": "r"}},
    "wipe":    {"name": "擦除",     "duration_ms": 400, "element": "wipe",    "attrs": {"dir": "r"}},
    "split":   {"name": "拆分",     "duration_ms": 500, "element": "split",   "attrs": {"orient": "horz", "dir": "out"}},
    "strips":  {"name": "条纹",     "duration_ms": 500, "element": "strips",  "attrs": {"dir": "rd"}},
    "cover":   {"name": "覆盖",     "duration_ms": 500, "element": "cover",   "attrs": {"dir": "r"}},
    "random":  {"name": "随机",     "duration_ms": 500, "element": "random",  "attrs": {}},
    "none":    {"name": "无转场",   "duration_ms": 0,   "element": None,      "attrs": {}},
}

ENTRANCE_ANIMATIONS = {
    "appear":       {"name": "出现",     "filter": None,                    "presetID": 1,  "presetSub": 0},
    "fade":         {"name": "淡入",     "filter": "fade",                  "presetID": 10, "presetSub": 0},
    "fly":          {"name": "飞入",     "filter": "slide(fromBottom)",     "presetID": 2,  "presetSub": 4},
    "cut":          {"name": "切入",     "filter": "slide(fromLeft)",       "presetID": 42, "presetSub": 8},
    "zoom":         {"name": "缩放",     "filter": "image",                 "presetID": 23, "presetSub": 0},
    "wipe":         {"name": "擦除",     "filter": "wipe(left)",            "presetID": 22, "presetSub": 1},
    "split":        {"name": "拆分",     "filter": "barn(inVertical)",      "presetID": 16, "presetSub": 21},
    "blinds":       {"name": "百叶窗",   "filter": "blinds(horizontal)",    "presetID": 3,  "presetSub": 10},
    "checkerboard": {"name": "棋盘格",   "filter": "checkerboard(across)",  "presetID": 5,  "presetSub": 6},
    "dissolve":     {"name": "溶解",     "filter": "dissolve",              "presetID": 9,  "presetSub": 0},
    "random_bars":  {"name": "随机线条", "filter": "randombar(horizontal)", "presetID": 14, "presetSub": 10},
    "peek":         {"name": "窥视",     "filter": "wipe(down)",            "presetID": 12, "presetSub": 4},
    "wheel":        {"name": "轮子",     "filter": "wheel(4)",              "presetID": 21, "presetSub": 0},
    "box":          {"name": "盒状",     "filter": "box(in)",               "presetID": 4,  "presetSub": 0},
    "circle":       {"name": "圆形",     "filter": "circle(in)",            "presetID": 6,  "presetSub": 0},
    "diamond":      {"name": "菱形",     "filter": "diamond(in)",           "presetID": 8,  "presetSub": 0},
    "plus":         {"name": "十字",     "filter": "plus(in)",              "presetID": 13, "presetSub": 0},
    "strips":       {"name": "条纹",     "filter": "strips(downRight)",     "presetID": 18, "presetSub": 12},
    "wedge":        {"name": "楔形",     "filter": "wedge",                 "presetID": 20, "presetSub": 0},
    "stretch":      {"name": "伸展",     "filter": "stretch(across)",       "presetID": 17, "presetSub": 0},
    "expand":       {"name": "展开",     "filter": "stretch(across)",       "presetID": 50, "presetSub": 0},
    "swivel":       {"name": "旋转",     "filter": "wheel(1)",              "presetID": 19, "presetSub": 0},
}

ENTRANCE_MIXED_POOL = [
    "blinds", "checkerboard", "dissolve", "fly", "cut",
    "random_bars", "box", "split", "strips", "wedge", "wheel",
    "wipe", "expand", "fade", "swivel", "zoom",
]


def _pick_effect(mode, idx, offset=0):
    """解析动画模式 → 具体效果名"""
    if mode in ENTRANCE_ANIMATIONS: return mode
    if mode == "mixed":
        return "fade" if idx == 0 else ENTRANCE_MIXED_POOL[(idx - 1 + offset) % len(ENTRANCE_MIXED_POOL)]
    if mode == "random":
        import random
        return random.choice(ENTRANCE_MIXED_POOL)
    return "fade"


def _build_effect_xml(anim_name, shape_id, dur_ms, set_id, eff_id):
    """单个元素的动画效果 XML"""
    info = ENTRANCE_ANIMATIONS.get(anim_name, ENTRANCE_ANIMATIONS["fade"])
    set_block = (
        f'<p:set><p:cBhvr><p:cTn id="{set_id}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
    )
    if info["filter"]:
        return (set_block +
                f'<p:animEffect transition="in" filter="{info["filter"]}">'
                f'<p:cBhvr><p:cTn id="{eff_id}" dur="{dur_ms}"/>'
                f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
                f'</p:cBhvr></p:animEffect>')
    return set_block


def _build_sequence_timing_xml(targets, duration=0.4, stagger=0.5, trigger="after-previous"):
    """多元素序列动画 timing XML（与 ppt-master create_sequence_timing_xml 一致）"""
    if not targets: return ""
    dur_ms = int(duration * 1000)
    stagger_ms = int(stagger * 1000)
    next_id = 3

    # after-previous 模式（默认，自动播放）
    elapsed_ms = 0
    inner_steps = []
    for i, (shape_id, delay_ms, anim_name) in enumerate(targets):
        if i > 0: elapsed_ms += dur_ms + delay_ms
        wrapper_id = next_id; leaf_id = next_id + 1
        set_id = next_id + 2; eff_id = next_id + 3
        next_id += 4
        effect = _build_effect_xml(anim_name, shape_id, dur_ms, set_id, eff_id)
        preset_id = ENTRANCE_ANIMATIONS.get(anim_name, ENTRANCE_ANIMATIONS["fade"])["presetID"]
        preset_sub = ENTRANCE_ANIMATIONS.get(anim_name, ENTRANCE_ANIMATIONS["fade"])["presetSub"]
        inner_steps.append(
            f'<p:par><p:cTn id="{wrapper_id}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{elapsed_ms}"/></p:stCondLst>'
            f'<p:childTnLst><p:par>'
            f'<p:cTn id="{leaf_id}" presetID="{preset_id}" presetClass="entr" '
            f'presetSubtype="{preset_sub}" fill="hold" nodeType="afterEffect">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{effect}</p:childTnLst>'
            f'</p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
        )

    outer_id = next_id
    all_steps = (
        f'<p:par><p:cTn id="{outer_id}" fill="hold">'
        f'<p:stCondLst><p:cond delay="indefinite"/>'
        f'<p:cond evt="onBegin" delay="0"><p:tn val="2"/></p:cond></p:stCondLst>'
        f'<p:childTnLst>{"".join(inner_steps)}</p:childTnLst>'
        f'</p:cTn></p:par>'
    )

    bld_list = "".join(f'<p:bldP spid="{sid}" grpId="0"/>' for sid, _, _ in targets)

    return (
        f'<p:timing xmlns:p="{NS_P}" xmlns:a="{NS_A}">'
        f'<p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f'{all_steps}'
        f'</p:childTnLst></p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
        f'<p:bldLst>{bld_list}</p:bldLst>'
        f'</p:timing>'
    )


def _build_transition_xml(effect="fade", duration_ms=500):
    """转场 XML"""
    tr = TRANSITIONS.get(effect)
    if not tr or not tr.get("element"): return ""
    elem = tr["element"]
    attrs = " ".join(f'{k}="{v}"' for k, v in tr.get("attrs", {}).items())
    if attrs: attrs = " " + attrs
    ns14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    return (f'<p:transition xmlns:p="{NS_P}" xmlns:p14="{ns14}" '
            f'p14:dur="{duration_ms}" advClick="1"><p:{elem}{attrs}/></p:transition>')


# ============================================================
#  公开 API
# ============================================================

def _repair_svg(svg_str):
    """修复 DeepSeek 生成的常见 SVG 语法问题"""
    # 修复 & 符号
    svg_str = svg_str.replace("&", "&amp;").replace("&amp;amp;", "&amp;").replace("&amp;#", "&#")

    # 确保有 </svg> 结尾
    if "</svg>" not in svg_str:
        svg_str += "</svg>"

    # 修复未关闭的标签（内层先关闭：tspan → text → g）
    for tag in ["tspan", "text", "defs", "g"]:
        opened = len(re.findall(rf"<{tag}[\s>]", svg_str))
        closed = svg_str.count(f"</{tag}>")
        if opened > closed:
            svg_str = svg_str.replace("</svg>", f"</{tag}>" * (opened - closed) + "</svg>")

    return svg_str


def svg_to_slide(slide, svg_str):
    """转换 SVG → slide 形状，返回 anim_targets"""
    svg_str = _repair_svg(svg_str)
    converter = SvgConverter()
    return converter.convert(slide, svg_str)


def svgs_to_pptx(pages, output_path, animation="mixed", transition="fade",
                 stagger=0.5, duration=0.4):
    """将多个页面转换为 PPTX

    AI 负责：SVG 设计 + 演讲稿
    代码负责：动画 + 转场（规则自动分配，不让 AI 决定）

    Args:
        pages: 页面列表，每项可以是：
            - 纯 SVG 字符串
            - dict: {"svg": "...", "notes": "演讲稿"}
        output_path: 输出文件路径
        animation: 入场动画模式（mixed/fade/none）
        transition: 转场效果（fade/push/wipe/none）
        stagger: 元素间动画间隔（秒）
        duration: 单个动画持续时间（秒）
    """
    global _shape_id_counter
    _shape_id_counter = 100

    prs = Presentation()
    prs.slide_width = Emu(px(1280))
    prs.slide_height = Emu(px(720))
    blank = prs.slide_layouts[6]
    mixed_offset = 0
    total = len(pages)

    # 转场规则：首页 fade，中间页 push，末页 fade
    def _auto_transition(idx):
        if idx == 0 or idx == total - 1: return "fade"
        return transition

    for i, page in enumerate(pages):
        if isinstance(page, dict):
            svg_str = page.get("svg", "")
            notes = page.get("notes", "")
        else:
            svg_str = page
            notes = ""

        slide = prs.slides.add_slide(blank)
        anim_targets = []

        try:
            anim_targets = svg_to_slide(slide, svg_str)
        except Exception as e:
            sp_tree = slide._element.find(f".//{qn('p:spTree')}")
            err_xml = _text_shape_xml(_next_id(), px(60), px(300), px(1160), px(100),
                                      _run_xml(f"Page {i+1} error: {str(e)[:80]}", 16), "l")
            sp_tree.append(etree.fromstring(err_xml))

        sld = slide._element

        # 1. 演讲稿 → speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        # 2. 转场（代码规则，首尾 fade，中间用指定效果）
        tr_name = _auto_transition(i)
        if tr_name and tr_name != "none":
            tr_info = TRANSITIONS.get(tr_name, TRANSITIONS.get("fade"))
            if tr_info:
                tr_xml = _build_transition_xml(tr_name, tr_info["duration_ms"])
                if tr_xml:
                    sld.append(etree.fromstring(tr_xml))

        # 3. 入场动画（代码规则，mixed 模式自动分配）
        if animation and animation != "none" and anim_targets:
            stagger_ms = int(stagger * 1000)
            seq_targets = []
            for idx, (sid, svg_id) in enumerate(anim_targets):
                effect = _pick_effect(animation, idx, mixed_offset)
                delay = 0 if idx == 0 else stagger_ms
                seq_targets.append((sid, delay, effect))
            mixed_offset += max(0, len(anim_targets) - 1)

            timing_xml = _build_sequence_timing_xml(seq_targets, duration, stagger)
            if timing_xml:
                sld.append(etree.fromstring(timing_xml))

    prs.save(output_path)
    return len(pages)
