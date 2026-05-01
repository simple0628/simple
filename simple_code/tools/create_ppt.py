"""创建 PowerPoint 演示文稿（多框架 + 转场 + 入场动画）"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
FONT = "Microsoft YaHei"

def _px(v): return Emu(int(v * 12192000 / 1280))
def _py(v): return Emu(int(v * 6858000 / 720))


definition = {
    "type": "function",
    "function": {
        "name": "create_ppt",
        "description": """创建 PowerPoint 演示文稿（.pptx）。使用前先调用 ppt_frameworks 工具查看可用框架。

参数说明：
- path（必填）：输出文件路径，.pptx 结尾
- framework（推荐）：设计框架名，先调用 ppt_frameworks 查看
  可选值：mckinsey / google / academic / government / pixel_retro / ai_enterprise / anthropic / smart_red / tech_blue / exhibit
- slides（必填）：幻灯片数组，每个元素包含：
  - layout: "title"/"content"/"section"
  - title: 标题
  - subtitle: 副标题（title/section 用）
  - bullets: 项目符号列表（content 用）
  - transition: 转场效果（fade/push/wipe/split/strips/cover/random/none）
  - entrance: 入场动画，先调用 ppt_frameworks 查看全部 22 种
    常用：fade/fly/zoom/wipe/blinds/dissolve/box/circle/diamond
    特殊：mixed(首元素淡入其余循环变化) / random(随机)
  - notes: 演讲者备注""",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "framework": {"type": "string"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout": {"type": "string", "enum": ["title", "content", "section"]},
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                            "transition": {"type": "string"},
                            "entrance": {"type": "string"},
                            "notes": {"type": "string"},
                        },
                        "required": ["layout", "title"],
                    },
                },
            },
            "required": ["path", "slides"],
        },
    },
}

needs_pause = False

def label(args):
    return f"正在设计 PPT 结构: {args.get('path', '')}"


# ============================================================
#  基础绘图
# ============================================================

def _hex(s):
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def _rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _rect_a(slide, l, t, w, h, color, alpha):
    shape = _rect(slide, l, t, w, h, color)
    s = shape._element.spPr.find(f".//{qn('a:srgbClr')}")
    if s is not None: etree.SubElement(s, qn("a:alpha")).set("val", str(alpha * 1000))
    return shape

def _outline(slide, l, t, w, h, color, alpha=100):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.background(); shape.line.color.rgb = color; shape.line.width = Pt(1)
    if alpha < 100:
        ln = shape._element.spPr.find(qn("a:ln"))
        if ln is not None:
            s = ln.find(f".//{qn('a:srgbClr')}")
            if s is not None: etree.SubElement(s, qn("a:alpha")).set("val", str(alpha * 1000))
    return shape

def _grad(slide, l, t, w, h, c1, c2, angle=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.line.fill.background()
    f = shape.fill; f.gradient(); f.gradient_angle = angle
    stops = f.gradient_stops
    stops[0].color.rgb = c1; stops[0].position = 0.0
    stops[1].color.rgb = c2; stops[1].position = 1.0
    return shape

def _oval(slide, l, t, w, h, color, alpha=100):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha < 100:
        s = shape._element.spPr.find(f".//{qn('a:srgbClr')}")
        if s is not None: etree.SubElement(s, qn("a:alpha")).set("val", str(alpha * 1000))
    return shape

def _tri(slide, l, t, w, h, color, alpha=100):
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, l, t, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha < 100:
        s = shape._element.spPr.find(f".//{qn('a:srgbClr')}")
        if s is not None: etree.SubElement(s, qn("a:alpha")).set("val", str(alpha * 1000))
    return shape

def _text(slide, text, l, t, w, h, size=18, color=None, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.name = FONT
    p.font.color.rgb = color or RGBColor(0x33,0x33,0x33)
    p.font.bold = bold; p.alignment = align
    p.line_spacing = Pt(int(size * 1.5))
    return box

def _set_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _bullets(slide, items, l, t, w, h, color, accent_hex):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(16); p.font.name = FONT
        p.font.color.rgb = color; p.space_after = Pt(14); p.line_spacing = Pt(30)
        pf = p._pPr
        if pf is None: pf = etree.SubElement(p._p, qn("a:pPr"))
        bn = pf.find(qn("a:buNone"))
        if bn is not None: pf.remove(bn)
        bc = pf.find(qn("a:buChar"))
        if bc is None: bc = etree.SubElement(pf, qn("a:buChar"))
        bc.set("char", "●")
        bcl = pf.find(qn("a:buClr"))
        if bcl is None: bcl = etree.SubElement(pf, qn("a:buClr"))
        else: bcl.clear()
        etree.SubElement(bcl, qn("a:srgbClr")).set("val", accent_hex.lstrip("#"))
        pf.set("marL", str(Emu(457200))); pf.set("indent", str(Emu(-228600)))
    return box

W = RGBColor(0xFF, 0xFF, 0xFF)


# ============================================================
#  装饰原语（各框架组合使用）
# ============================================================

def _deco_left_stripe(slide, c1, c2):
    """左侧渐变竖条 + 强调线"""
    _grad(slide, 0, 0, _px(12), SLIDE_H, c1, c2, 270)
    _rect(slide, _px(12), 0, _px(3), SLIDE_H, c2)

def _deco_geometric_right(slide, c):
    """右侧几何方块轮廓"""
    _rect_a(slide, _px(780), 0, _px(500), SLIDE_H, c, 8)
    _outline(slide, _px(850), _py(120), _px(360), _py(360), c, 15)
    _outline(slide, _px(900), _py(170), _px(260), _py(260), c, 10)
    _rect_a(slide, _px(950), _py(220), _px(160), _py(160), c, 3)

def _deco_four_color_bar(slide, colors):
    """顶部四色条"""
    seg = SLIDE_W // len(colors)
    for i, c in enumerate(colors):
        _rect(slide, Emu(seg * i), 0, Emu(seg), Pt(5), c)

def _deco_grid_lines(slide, alpha=3):
    """背景网格线"""
    for y in [180, 360, 540]:
        _rect_a(slide, 0, _py(y), SLIDE_W, Pt(1), W, alpha)
    for x in [320, 640, 960]:
        _rect_a(slide, _px(x), 0, Pt(1), SLIDE_H, W, alpha)

def _deco_center_divider(slide, c_line, c_dot, y=340):
    """居中金色圆点 + 对称线"""
    _rect_a(slide, _px(500), _py(y), _px(120), _py(2), W, 40)
    _oval(slide, _px(634), _py(y - 6), _px(12), _py(12), c_dot)
    _rect_a(slide, _px(660), _py(y), _px(120), _py(2), W, 40)

def _deco_insight_bar(slide, accent):
    """KEY INSIGHT 洞察栏"""
    _rect(slide, _px(60), _py(100), _px(1160), _py(50), _hex("#FFF8E1"))
    _rect(slide, _px(60), _py(100), _px(5), _py(50), accent)
    _text(slide, "KEY INSIGHT", _px(85), _py(108), _px(120), _py(30),
          size=11, color=accent, bold=True)

def _deco_footer(slide, light_bg, line_c, primary, pn, total):
    """页脚：灰色底 + 线 + 蓝底页码"""
    _rect(slide, 0, _py(640), SLIDE_W, _py(80), light_bg)
    _rect(slide, 0, _py(640), SLIDE_W, Pt(1), line_c)
    _rect(slide, _px(1140), _py(655), _px(80), _py(30), primary)
    _text(slide, f"{pn}/{total}", _px(1140), _py(655), _px(80), _py(30),
          size=11, color=W, bold=True, align=PP_ALIGN.CENTER)

def _deco_triangles(slide, c1, c2, alpha1=12, alpha2=8):
    """三角形装饰"""
    _tri(slide, _px(0), _py(0), _px(300), _py(250), c1, alpha1)
    _tri(slide, _px(980), _py(470), _px(300), _py(250), c2, alpha2)

def _deco_neon_lines(slide, c):
    """霓虹双线"""
    _rect(slide, 0, _py(2), SLIDE_W, _py(2), c)
    _rect(slide, 0, _py(7), SLIDE_W, _py(1), c)

def _deco_pixel_corners(slide, c, alpha=20):
    """像素角落方块"""
    s = _px(15)
    for x, y in [(30, 30), (1235, 30), (30, 675), (1235, 675)]:
        _rect_a(slide, _px(x), _py(y), s, s, c, alpha)
        _rect_a(slide, _px(x) + s, _py(y), s, s, c, alpha // 2)
        _rect_a(slide, _px(x), _py(y) + s, s, s, c, alpha // 2)


# ============================================================
#  框架布局实现
# ============================================================

def _resolve_theme(name):
    from simple_code.tools.ppt_frameworks import FRAMEWORKS
    default = FRAMEWORKS["mckinsey"]["colors"]
    if name is None: return default, "mckinsey"
    if isinstance(name, str) and name in FRAMEWORKS:
        return FRAMEWORKS[name]["colors"], name
    if isinstance(name, dict):
        base = dict(default); base.update(name)
        return base, "mckinsey"
    return default, "mckinsey"


# 框架分类映射：21 个框架 → 9 种视觉风格
_STYLE_MAP = {
    # A: 白底左对齐 + 左侧竖条 + 右侧几何轮廓 + 信息卡片
    "mckinsey": "A", "mckinsey_default": "A",
    # B: 白底居中 + 四色系统
    "google": "B",
    # C: 白底居中 + 蓝色顶条 + 红/橙左竖条 + 圆点分割线
    "academic": "C", "medical": "C", "cqu": "C",
    # D: 深色渐变底 + 科技网格/圆环 + 亮色强调条
    "government_blue": "D", "government_red": "D", "powerchina": "D", "powerchina_modern": "D",
    "catarc_business": "D", "catarc_standard": "D", "catarc_modern": "D",
    # E: 暗色底 + 霓虹线 + 像素方块
    "pixel_retro": "E",
    # F: 暗色底 + 网格线 + 强调色竖线
    "exhibit": "F", "anthropic": "F",
    # G: 浅灰底 + 三角形几何切割
    "smart_red": "G",
    # H: 白底 + 红蓝双色竖条 + 编号徽章 + 模块化
    "ai_enterprise": "H", "china_telecom": "H",
    # I: 白底左对齐 + 波浪/渐变 + 简洁
    "tech_blue": "I", "cmb": "I", "psychology": "I",
}

def _get_style(fw):
    return _STYLE_MAP.get(fw, "A")


def _build_title(slide, title, subtitle, t, fw, pn, total):
    p1, p2, acc = _hex(t["primary"]), _hex(t.get("secondary", t["primary"])), _hex(t.get("accent", t["primary"]))
    body_c, title_c = _hex(t.get("body_text", "#5D6D7E")), _hex(t.get("title_text", "#2C3E50"))
    light = _hex(t.get("light_bg", "#FAFBFC"))
    line_c = _hex(t.get("line", "#E5E7EB"))
    sty = _get_style(fw)

    if sty == "E":  # 像素风
        _set_bg(slide, _hex(t.get("light_bg", "#0D1117")))
        _deco_neon_lines(slide, p1)
        _deco_pixel_corners(slide, p1)
        _text(slide, title, _px(100), _py(250), _px(1080), _py(120),
              size=44, color=_hex(t["title_text"]), bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(50),
                  size=18, color=_hex(t["body_text"]), align=PP_ALIGN.CENTER)
        _rect(slide, 0, _py(716), SLIDE_W, _py(4), p1)

    elif sty == "F":  # 暗色咨询/AI
        _set_bg(slide, _hex(t.get("secondary", t["primary"])))
        _deco_grid_lines(slide, 3)
        _rect(slide, _px(60), _py(280), _px(5), _py(160), acc)
        _text(slide, title, _px(80), _py(250), _px(700), _py(120), size=44, color=W, bold=True)
        if subtitle:
            _text(slide, subtitle, _px(80), _py(400), _px(700), _py(50), size=18, color=W)
        _rect(slide, 0, _py(716), SLIDE_W, Pt(4), acc)

    elif sty == "B":  # 四色
        _set_bg(slide, W)
        bar_colors = [_hex(t["primary"]), _hex(t.get("red", t["primary"])),
                      _hex(t.get("accent", t["primary"])), _hex(t.get("green", t["primary"]))]
        seg_h = SLIDE_H // 4
        for i, c in enumerate(bar_colors):
            _rect(slide, 0, Emu(seg_h * i), _px(10), Emu(seg_h), c)
        _text(slide, title, _px(100), _py(260), _px(1080), _py(120),
              size=44, color=title_c, bold=True, align=PP_ALIGN.CENTER)
        _rect(slide, _px(340), _py(390), _px(150), Pt(4), bar_colors[0])
        _rect(slide, _px(490), _py(390), _px(70), Pt(4), bar_colors[1])
        _rect(slide, _px(560), _py(390), _px(70), Pt(4), bar_colors[2])
        _rect(slide, _px(630), _py(390), _px(170), Pt(4), bar_colors[3])
        if subtitle:
            _text(slide, subtitle, _px(100), _py(410), _px(1080), _py(50),
                  size=18, color=body_c, align=PP_ALIGN.CENTER)

    elif sty == "G":  # 三角几何
        _set_bg(slide, light)
        _tri(slide, _px(0), _py(0), _px(350), _py(300), p1, 15)
        _tri(slide, _px(930), _py(420), _px(350), _py(300), p1, 10)
        _tri(slide, _px(0), _py(420), _px(250), _py(300), _hex(t.get("accent", "#333333")), 12)
        _text(slide, title, _px(100), _py(260), _px(1080), _py(120),
              size=44, color=title_c, bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(50),
                  size=18, color=body_c, align=PP_ALIGN.CENTER)

    elif sty == "C":  # 学术/医学
        _set_bg(slide, W)
        _rect(slide, 0, 0, SLIDE_W, _py(6), p1)
        _rect(slide, 0, 0, _px(6), SLIDE_H, acc)
        _text(slide, title, _px(60), _py(250), _px(1100), _py(120),
              size=44, color=title_c, bold=True, align=PP_ALIGN.CENTER)
        _oval(slide, _px(620), _py(390), _px(10), _py(10), p1)
        _rect(slide, _px(450), _py(394), _px(160), _py(2), line_c)
        _rect(slide, _px(640), _py(394), _px(160), _py(2), line_c)
        if subtitle:
            _text(slide, subtitle, _px(60), _py(410), _px(1100), _py(50),
                  size=18, color=body_c, align=PP_ALIGN.CENTER)
        _rect(slide, 0, _py(680), SLIDE_W, Pt(1), line_c)

    elif sty == "D":  # 政务/工程深色
        _grad(slide, 0, 0, SLIDE_W, SLIDE_H, p1, _hex(t.get("secondary", t["primary"])), 315)
        _deco_grid_lines(slide, 2)
        _rect(slide, _px(60), _py(280), _px(6), _py(160), acc)
        _text(slide, title, _px(80), _py(260), _px(700), _py(120), size=44, color=W, bold=True)
        if subtitle:
            _text(slide, subtitle, _px(80), _py(400), _px(700), _py(50), size=18, color=W)
        _rect(slide, 0, _py(716), SLIDE_W, Pt(4), acc)

    elif sty == "H":  # 红蓝双色模块化
        _set_bg(slide, W)
        _rect(slide, 0, 0, _px(30), SLIDE_H // 2, p1)
        _rect(slide, 0, SLIDE_H // 2, _px(30), SLIDE_H // 2, p2)
        _rect(slide, _px(30), 0, _px(30), SLIDE_H, acc)
        _text(slide, title, _px(100), _py(260), _px(1080), _py(120),
              size=44, color=title_c, bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(50),
                  size=18, color=body_c, align=PP_ALIGN.CENTER)
        _rect(slide, 0, _py(700), SLIDE_W, _py(20), p2)

    elif sty == "I":  # 简洁渐变波浪
        _set_bg(slide, W)
        _grad(slide, 0, 0, _px(500), SLIDE_H, p1, p2, 315)
        _text(slide, title, _px(60), _py(260), _px(420), _py(120), size=40, color=W, bold=True)
        if subtitle:
            _text(slide, subtitle, _px(60), _py(400), _px(420), _py(50), size=18, color=W)
        _rect(slide, 0, _py(680), SLIDE_W, _py(6), acc)
        _rect(slide, 0, _py(690), SLIDE_W, _py(3), p1)

    else:  # A: 默认麦肯锡
        _set_bg(slide, W)
        _deco_left_stripe(slide, p1, acc)
        _deco_geometric_right(slide, p1)
        _rect(slide, _px(60), _py(55), _px(180), _py(5), p1)
        _rect_a(slide, _px(60), _py(65), _px(100), _py(3), p2, 50)
        _text(slide, title, _px(60), _py(250), _px(700), _py(120), size=44, color=title_c, bold=True)
        if subtitle:
            _text(slide, subtitle, _px(60), _py(390), _px(700), _py(50), size=20, color=body_c)
        _rect(slide, _px(60), _py(450), _px(440), _py(2), line_c)
        _rect(slide, _px(60), _py(456), _px(140), _py(2), p1)
        _rect(slide, _px(60), _py(520), _px(400), _py(100), light)
        _rect(slide, _px(60), _py(520), _px(4), _py(100), p1)
        _rect(slide, 0, _py(718), SLIDE_W, Pt(6), acc)


def _build_section(slide, title, subtitle, t, fw, pn, total):
    p1, p2, acc = _hex(t["primary"]), _hex(t.get("secondary", t["primary"])), _hex(t.get("accent", t["primary"]))
    sty = _get_style(fw)

    if sty == "E":
        _set_bg(slide, _hex(t.get("light_bg", "#0D1117")))
        _deco_neon_lines(slide, p1); _deco_pixel_corners(slide, p1)
        _text(slide, title, _px(100), _py(300), _px(1080), _py(80),
              size=40, color=_hex(t["title_text"]), bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(40),
                  size=18, color=_hex(t["body_text"]), align=PP_ALIGN.CENTER)

    elif sty == "B":
        bar_colors = [_hex(t["primary"]), _hex(t.get("red", t["primary"])),
                      _hex(t.get("accent", t["primary"])), _hex(t.get("green", t["primary"]))]
        _grad(slide, 0, 0, SLIDE_W, SLIDE_H, _hex(t.get("secondary", "#1A237E")), p1, 315)
        _deco_four_color_bar(slide, bar_colors)
        _text(slide, title, _px(100), _py(300), _px(1080), _py(80),
              size=40, color=W, bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(40),
                  size=18, color=W, align=PP_ALIGN.CENTER)

    elif sty == "G":
        _set_bg(slide, _hex(t.get("light_bg", "#F5F5F7")))
        _tri(slide, _px(0), _py(0), _px(250), _py(200), p1, 12)
        _tri(slide, _px(1030), _py(520), _px(250), _py(200), p1, 8)
        _text(slide, title, _px(100), _py(300), _px(1080), _py(80),
              size=40, color=_hex(t["title_text"]), bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(40),
                  size=18, color=_hex(t["body_text"]), align=PP_ALIGN.CENTER)

    elif sty == "C":
        _set_bg(slide, p1)
        _rect(slide, 0, 0, _px(6), SLIDE_H, acc)
        _rect_a(slide, _px(800), _py(100), _px(350), _py(350), W, 5)
        _text(slide, title, _px(80), _py(300), _px(700), _py(80), size=40, color=W, bold=True)
        if subtitle:
            _text(slide, subtitle, _px(80), _py(400), _px(700), _py(40), size=18, color=W)
        _text(slide, f"{pn}/{total}", _px(60), _py(700), _px(100), _py(20), size=11, color=W)

    elif sty == "H":
        _set_bg(slide, W)
        _rect(slide, 0, 0, _px(30), SLIDE_H, p1)
        _rect(slide, _px(1250), 0, _px(30), SLIDE_H, p2)
        _rect_a(slide, _px(100), _py(200), _px(100), _py(100), p1, 8)
        _text(slide, title, _px(100), _py(300), _px(1080), _py(80),
              size=40, color=_hex(t["title_text"]), bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(40),
                  size=18, color=_hex(t["body_text"]), align=PP_ALIGN.CENTER)
        _rect(slide, _px(500), _py(450), _px(280), Pt(3), p1)
        _rect(slide, _px(500), _py(456), _px(280), Pt(2), p2)

    elif sty == "I":
        _grad(slide, 0, 0, SLIDE_W, SLIDE_H, p1, p2, 315)
        _text(slide, title, _px(100), _py(300), _px(1080), _py(80),
              size=40, color=W, bold=True, align=PP_ALIGN.CENTER)
        if subtitle:
            _text(slide, subtitle, _px(100), _py(400), _px(1080), _py(40),
                  size=18, color=W, align=PP_ALIGN.CENTER)
        _rect(slide, 0, _py(700), SLIDE_W, Pt(4), acc)
        _text(slide, f"{pn}/{total}", _px(60), _py(700), _px(100), _py(20), size=11, color=W)

    else:  # A, D, F
        dark = _hex(t.get("secondary", t["primary"]))
        _grad(slide, 0, 0, SLIDE_W, SLIDE_H, p1, dark, 315)
        _deco_grid_lines(slide)
        _rect(slide, _px(60), _py(280), _px(5), _py(160), acc)
        _text(slide, title, _px(80), _py(300), _px(600), _py(80), size=40, color=W, bold=True)
        _deco_center_divider(slide, W, acc)
        if subtitle:
            _text(slide, subtitle, _px(80), _py(400), _px(600), _py(40), size=18, color=W)
        _rect_a(slide, 0, _py(680), SLIDE_W, _py(40), W, 5)
        _text(slide, f"{pn}/{total}", _px(60), _py(700), _px(100), _py(20), size=11, color=W)


def _build_content(slide, title, bullet_items, t, fw, pn, total):
    p1, p2, acc = _hex(t["primary"]), _hex(t.get("secondary", t["primary"])), _hex(t.get("accent", t["primary"]))
    body_c = _hex(t.get("body_text", "#5D6D7E"))
    title_c = _hex(t.get("title_text", "#2C3E50"))
    light = _hex(t.get("light_bg", "#FAFBFC"))
    line_c = _hex(t.get("line", "#E5E7EB"))
    sty = _get_style(fw)
    acc_hex = t.get("accent", t["primary"])

    if sty == "E":
        _set_bg(slide, _hex(t.get("light_bg", "#0D1117")))
        _deco_neon_lines(slide, p1); _deco_pixel_corners(slide, p1, 15)
        _text(slide, title, _px(80), _py(30), _px(1100), _py(40), size=22, color=p1, bold=True)
        _rect(slide, _px(80), _py(75), _px(200), _py(2), p1)
        if bullet_items:
            _bullets(slide, bullet_items, _px(100), _py(100), _px(1080), _py(520),
                     _hex(t["body_text"]), acc_hex)
        _text(slide, f"{pn}/{total}", _px(1160), _py(690), _px(80), _py(20),
              size=11, color=_hex(t["body_text"]), align=PP_ALIGN.RIGHT)

    elif sty == "B":
        _set_bg(slide, W)
        bar_colors = [_hex(t["primary"]), _hex(t.get("red", t["primary"])),
                      _hex(t.get("accent", t["primary"])), _hex(t.get("green", t["primary"]))]
        _deco_four_color_bar(slide, bar_colors)
        _rect(slide, 0, Pt(5), SLIDE_W, _py(75), light)
        _rect(slide, 0, _py(80), SLIDE_W, Pt(1), line_c)
        _text(slide, title, _px(60), _py(22), _px(900), _py(45), size=22, color=title_c, bold=True)
        _rect(slide, _px(60), _py(72), _px(80), Pt(3), bar_colors[0])
        _rect(slide, _px(140), _py(72), _px(40), Pt(3), bar_colors[1])
        _rect(slide, _px(180), _py(72), _px(40), Pt(3), bar_colors[2])
        _rect(slide, _px(220), _py(72), _px(80), Pt(3), bar_colors[3])
        if bullet_items:
            _bullets(slide, bullet_items, _px(80), _py(100), _px(1100), _py(520), body_c, acc_hex)
        _deco_footer(slide, light, line_c, _hex(t["primary"]), pn, total)

    elif sty == "G":
        _set_bg(slide, _hex(t.get("light_bg", "#F5F5F7")))
        _tri(slide, _px(1100), _py(0), _px(180), _py(80), p1, 10)
        _rect(slide, 0, 0, SLIDE_W, _py(80), W)
        _text(slide, title, _px(60), _py(22), _px(900), _py(45), size=22, color=title_c, bold=True)
        _rect(slide, _px(60), _py(72), _px(60), Pt(3), p1)
        _rect(slide, _px(125), _py(72), _px(30), Pt(3), _hex(t.get("secondary", "#F0964D")))
        if bullet_items:
            _bullets(slide, bullet_items, _px(80), _py(100), _px(1100), _py(520), body_c, acc_hex)
        _deco_footer(slide, light, line_c, p1, pn, total)

    elif sty == "C":
        _set_bg(slide, W)
        _rect(slide, 0, 0, SLIDE_W, _py(70), p1)
        _rect(slide, 0, 0, _px(6), _py(70), acc)
        _text(slide, title, _px(60), _py(15), _px(900), _py(40), size=22, color=W, bold=True)
        _rect(slide, _px(60), _py(90), _px(1160), _py(50), _hex(t.get("light_bg", "#E8F4FC")))
        _rect(slide, _px(60), _py(90), _px(5), _py(50), p1)
        if bullet_items:
            _bullets(slide, bullet_items, _px(85), _py(160), _px(1100), _py(460), body_c, acc_hex)
        _deco_footer(slide, _hex(t.get("light_bg", "#F5F7FA")), line_c, p1, pn, total)

    elif sty == "H":
        _set_bg(slide, W)
        _rect(slide, 0, 0, SLIDE_W, Pt(4), p1)
        _rect(slide, _px(60), _py(15), _px(8), _py(40), p1)
        _rect(slide, _px(60), _py(15), _px(40), _py(40), p1)
        _text(slide, title, _px(115), _py(15), _px(900), _py(40), size=22, color=title_c, bold=True)
        if bullet_items:
            _bullets(slide, bullet_items, _px(85), _py(80), _px(1100), _py(540), body_c, acc_hex)
        _rect(slide, 0, _py(640), SLIDE_W, _py(80), light)
        _rect(slide, 0, _py(640), SLIDE_W, Pt(1), line_c)
        _rect(slide, _px(60), _py(655), _px(4), _py(20), p1)
        _rect(slide, _px(1140), _py(655), _px(80), _py(30), p1)
        _text(slide, f"{pn}/{total}", _px(1140), _py(655), _px(80), _py(30),
              size=11, color=W, bold=True, align=PP_ALIGN.CENTER)

    elif sty == "I":
        _set_bg(slide, W)
        _rect(slide, 0, 0, SLIDE_W, Pt(4), p1)
        _rect(slide, _px(60), _py(15), _px(10), _py(40), p1)
        _text(slide, title, _px(85), _py(15), _px(900), _py(40), size=22, color=title_c, bold=True)
        _rect(slide, _px(85), _py(60), _px(200), Pt(2), acc)
        if bullet_items:
            _bullets(slide, bullet_items, _px(85), _py(85), _px(1100), _py(530), body_c, acc_hex)
        _deco_footer(slide, light, line_c, p1, pn, total)

    else:  # A, D, F 默认
        _set_bg(slide, W)
        _grad(slide, 0, 0, SLIDE_W, _py(5), p1, p2)
        _rect(slide, 0, _py(5), SLIDE_W, _py(75), light)
        _rect(slide, 0, _py(80), SLIDE_W, Pt(1), line_c)
        _rect(slide, _px(60), _py(22), _px(50), _py(40), p1)
        _text(slide, title, _px(130), _py(22), _px(900), _py(45), size=22, color=title_c, bold=True)
        _deco_insight_bar(slide, acc)
        _rect(slide, _px(60), _py(170), _px(4), _py(450), p1)
        if bullet_items:
            _bullets(slide, bullet_items, _px(85), _py(170), _px(1100), _py(450), body_c, acc_hex)
        _deco_footer(slide, light, line_c, p1, pn, total)


# ============================================================
#  转场
# ============================================================

def _add_transition(slide, effect, dur_ms):
    """添加转场效果（支持全部 8 种）"""
    from simple_code.tools.ppt_frameworks import TRANSITIONS
    if not effect or effect == "none":
        return
    tr = TRANSITIONS.get(effect)
    if not tr or not tr.get("element"):
        return
    elem = tr["element"]
    attrs = " ".join(f'{k}="{v}"' for k, v in tr.get("attrs", {}).items())
    if attrs:
        attrs = " " + attrs
    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    xml = f'<p:transition xmlns:p="{ns}" xmlns:p14="{ns14}" p14:dur="{dur_ms}" advClick="1"><p:{elem}{attrs}/></p:transition>'
    slide._element.append(etree.fromstring(xml))


# ============================================================
#  入场动画
# ============================================================

def _resolve_entrance(name, idx=0, slide_idx=0):
    """解析入场动画名 → 实际动画名（支持 mixed/random）"""
    from simple_code.tools.ppt_frameworks import ENTRANCE_ANIMATIONS, ENTRANCE_MIXED_POOL
    if name in ENTRANCE_ANIMATIONS:
        return name
    if name == "mixed":
        if idx == 0:
            return "fade"
        return ENTRANCE_MIXED_POOL[(idx - 1 + slide_idx) % len(ENTRANCE_MIXED_POOL)]
    if name == "random":
        import random
        return random.choice(ENTRANCE_MIXED_POOL)
    return None


def _add_entrance(slide, shape_id, anim_name="fade", dur_ms=800, delay_ms=0):
    """给指定 shape 添加入场动画（支持全部 22 种）"""
    from simple_code.tools.ppt_frameworks import ENTRANCE_ANIMATIONS
    if not anim_name or anim_name == "none" or anim_name not in ENTRANCE_ANIMATIONS:
        return
    info = ENTRANCE_ANIMATIONS[anim_name]

    ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    set_xml = (f'<p:set xmlns:p="{ns}">'
               f'<p:cBhvr><p:cTn id="5" dur="1" fill="hold">'
               f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
               f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
               f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
               f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>')

    if info["filter"]:
        effect_xml = (f'{set_xml}'
                      f'<p:animEffect xmlns:p="{ns}" '
                      f'transition="in" filter="{info["filter"]}">'
                      f'<p:cBhvr><p:cTn id="6" dur="{dur_ms}"/>'
                      f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
                      f'</p:cBhvr></p:animEffect>')
    else:
        effect_xml = set_xml

    timing_xml = (
        f'<p:timing xmlns:p="{ns}"><p:tnLst><p:par>'
        f'<p:cTn id="1" dur="indefinite" nodeType="tmRoot"><p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="none">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
        f'<p:par><p:cTn id="3" fill="hold">'
        f'<p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>'
        f'<p:childTnLst><p:par><p:cTn id="4" fill="hold"><p:childTnLst>'
        f'{effect_xml}'
        f'</p:childTnLst></p:cTn></p:par></p:childTnLst>'
        f'</p:cTn></p:par>'
        f'</p:childTnLst></p:cTn></p:seq>'
        f'</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>'
    )
    slide._element.append(etree.fromstring(timing_xml))


# ============================================================
#  主执行
# ============================================================

def execute(args, app=None, **kwargs):
    path = args["path"]
    slides_data = args["slides"]
    theme, fw_name = _resolve_theme(args.get("framework"))

    if app:
        app.finish_tool(success=True)
        answer = app.preview_ppt(slides_data)
        if answer == "(用户取消)":
            return "用户取消了 PPT 创建"
        if answer not in ("(用户未输入内容)", ""):
            return f"用户修改意见: {answer}\n请根据意见修改 slides 内容后重新调用 create_ppt。"
        app.write_tool(f"正在创建 PPT: {path}")

    from simple_code.tools.ppt_frameworks import TRANSITIONS, TRANSITION_DEFAULTS

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = len(slides_data)

    for idx, s in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        layout = s.get("layout", "content")
        title = s.get("title", "")
        subtitle = s.get("subtitle", "")
        bullet_items = s.get("bullets", [])
        notes = s.get("notes", "")
        pn = idx + 1

        if layout == "title":
            _build_title(slide, title, subtitle, theme, fw_name, pn, total)
        elif layout == "section":
            _build_section(slide, title, subtitle, theme, fw_name, pn, total)
        else:
            _build_content(slide, title, bullet_items, theme, fw_name, pn, total)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        # 转场
        tr_name = s.get("transition", TRANSITION_DEFAULTS.get(layout, "fade"))
        tr_info = TRANSITIONS.get(tr_name, TRANSITIONS["fade"])
        _add_transition(slide, tr_name, tr_info["duration_ms"])

        # 入场动画（应用到最后一个 shape，通常是内容区域）
        entrance_raw = s.get("entrance")
        if entrance_raw and entrance_raw != "none" and len(slide.shapes) > 0:
            resolved = _resolve_entrance(entrance_raw, idx=0, slide_idx=idx)
            if resolved:
                last_shape = slide.shapes[-1]
                _add_entrance(slide, last_shape.shape_id, resolved)

    prs.save(path)

    try:
        simple_dir = os.path.join(os.getcwd(), "simple")
        ppt_md = os.path.join(simple_dir, "simple-ppt.md")
        from datetime import date
        first_title = slides_data[0].get("title", "未知") if slides_data else "未知"
        line = f"- {date.today()} | {first_title} | {total}页 | {os.path.basename(path)}\n"
        if os.path.exists(simple_dir):
            with open(ppt_md, "a", encoding="utf-8") as f:
                f.write(line)
    except Exception:
        pass

    return f"PPT 已创建: {path}（共 {total} 页）"
