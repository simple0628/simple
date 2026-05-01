"""simple 入口：Textual TUI 主循环"""

import sys
import os
import re
import time
import threading

if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8")

from openai import OpenAI

from simple_code import __version__
from simple_code.config import (
    get_provider, set_provider_key,
    test_provider_key, load_skills, SKILLS_DIR
)
from simple_code.chat import build_system_prompt, chat_round, save_memory
from simple_code.ui import SimpleApp
from simple_code import state


def main():
    # 获取模型配置
    provider = get_provider()
    if not provider["api_key"]:
        # 首次使用，引导配置
        from simple_code.config import first_run_setup
        first_run_setup()
        provider = get_provider()

    client = OpenAI(api_key=provider["api_key"], base_url=provider["base_url"])
    model_name = provider["model"]

    cwd = os.getcwd()
    system_prompt = build_system_prompt(cwd, sys.platform)

    # 记忆文件夹
    simple_dir = os.path.join(cwd, "simple")
    os.makedirs(simple_dir, exist_ok=True)
    simple_md_path = os.path.join(simple_dir, "simple.md")
    simple_ppt_path = os.path.join(simple_dir, "simple-ppt.md")

    old_simple_md = os.path.join(cwd, "simple.md")
    if os.path.exists(old_simple_md) and not os.path.exists(simple_md_path):
        os.rename(old_simple_md, simple_md_path)

    if os.path.exists(simple_md_path):
        with open(simple_md_path, "r", encoding="utf-8") as f:
            system_prompt += f"\n## 项目记忆\n{f.read()}\n"

    if os.path.exists(simple_ppt_path):
        with open(simple_ppt_path, "r", encoding="utf-8") as f:
            system_prompt += f"\n## PPT 偏好记忆\n{f.read()}\n"

    messages = [{"role": "system", "content": system_prompt}]
    token_counter = {"total": 0, "round": 0}
    tool_logs = []

    def handle_submit(text):
        nonlocal client, model_name, provider

        skills = load_skills()
        cmd = text.strip()

        # --- 斜杠命令 ---
        if cmd == "/指南":
            app.write_system("")
            app.write_system("命令:")
            app.write_system("  /指南   显示帮助")
            app.write_system("  /模型   管理模型密钥")
            app.write_system("  /ppt    进入 PPT 制作模式")
            app.write_system("  /清空   清空对话历史")
            if skills:
                app.write_system("")
                app.write_system("自定义 Skill:")
                for name in sorted(skills):
                    desc = skills[name].split("\n")[0][:50]
                    app.write_system(f"  /{name}    {desc}")
            app.write_system("")
            app.write_system(f"  Skill 目录: {SKILLS_DIR}")
            app.write_system("")
            app.write_system("快捷键:")
            app.write_system("  Enter        发送消息")
            app.write_system("  ESC          中断当前任务")
            app.write_system("  Ctrl+C × 2   退出程序")
            app.write_system("  Ctrl+V       粘贴")
            app.write_system("  PageUp/Down  翻页浏览")
            app.write_system("")
            return

        if cmd == "/模型":
            _handle_model_command()
            return

        if cmd == "/ppt":
            from simple_code.ppt_mode import start as ppt_start
            ppt_start(app, client, model_name)
            return

        if cmd == "/重置":
            app.write_system("请使用 /模型 命令管理 API Key")
            return

        if cmd == "/清空":
            messages.clear()
            messages.append({"role": "system", "content": system_prompt})
            token_counter["total"] = 0
            tool_logs.clear()
            app.write_system("对话已清空")
            return

        # 自定义 skill
        user_input = text
        if cmd.startswith("/"):
            parts = cmd.split(None, 1)
            skill_name = parts[0][1:]
            extra = parts[1] if len(parts) > 1 else ""
            skills = load_skills()
            if skill_name in skills:
                skill_content = skills[skill_name]
                user_input = f"{skill_content}\n\n{extra}" if extra else skill_content
                app.write_system(f"已加载 skill: {skill_name}")
            else:
                app.write_warning(f"未知命令: /{skill_name}")
                app.write_system("输入 /指南 查看可用命令")
                return

        # --- 文件路径检测（拖拽文件自动处理）---
        user_input = _process_file_paths(user_input, app)

        # --- 正常对话 ---
        app.write_user(user_input)

        messages.append({"role": "user", "content": user_input})
        token_counter["round"] = 0
        tool_logs.clear()

        task_start_time = time.time()
        app.start_thinking()
        state.interrupt.clear()

        reply = chat_round(client, model_name, messages, app, tool_logs, token_counter, state.interrupt)

        was_interrupted = state.interrupt.is_set()
        state.interrupt.set()

        if was_interrupted:
            app.write_warning("已中断")

        app.stop_thinking()

        if save_memory(client, model_name, simple_md_path, user_input, reply, token_counter):
            app.write_system("记忆已保存")

        elapsed = int(time.time() - task_start_time)
        round_k = f"{token_counter['round'] / 1000:.1f}k"
        total_k = f"{token_counter['total'] / 1000:.1f}k"
        app.show_stats(elapsed, round_k, total_k)

    def _handle_model_command():
        nonlocal client, provider

        app.write_system("")
        app.write_system(f"当前模型: {provider['name']} ({provider['model']})")
        app.write_system(f"获取 API Key: {provider['key_url']}")
        app.write_system("")
        new_key = app.request_user_input("请输入新的 API Key（直接回车取消）:")
        if new_key in ("(用户未输入内容)", "(用户取消)"):
            return
        app.write_system("正在验证...")
        if test_provider_key(provider["id"], new_key):
            set_provider_key(provider["id"], new_key)
            client = OpenAI(api_key=new_key, base_url=provider["base_url"])
            app.write_system("API Key 已更新")
        else:
            app.write_warning("验证失败，请检查 API Key 是否正确")

    def _safe_handle_submit(text):
        try:
            handle_submit(text)
        except Exception as e:
            app.write_warning(f"出错: {e}")
            app.stop_thinking()

    def submit_in_thread(text):
        threading.Thread(target=_safe_handle_submit, args=(text,), daemon=True).start()

    app = SimpleApp(on_submit=submit_in_thread)
    app.version = __version__
    app.cwd = cwd
    app.has_memory = os.path.exists(simple_md_path)
    app.provider_name = provider["name"]
    app.model_name = model_name
    try:
        app.run()
    except KeyboardInterrupt:
        pass
    finally:
        if sys.platform == "win32":
            import ctypes
            kernel32 = ctypes.windll.kernel32
            stdin_handle = kernel32.GetStdHandle(-10)
            kernel32.SetConsoleMode(stdin_handle, 0x0007)


def _process_file_paths(user_input, app):
    """检测输入中的文件路径（拖拽产生），自动读取内容附加到消息中。"""
    # 提取带引号的路径或整行路径
    path_pattern = r'"([^"]+)"|\'([^\']+)\''
    matches = re.findall(path_pattern, user_input)
    candidates = [m[0] or m[1] for m in matches]

    # 如果没有引号包裹的路径，尝试整行作为路径
    if not candidates:
        stripped = user_input.strip().strip('"').strip("'")
        if os.path.exists(stripped):
            candidates = [stripped]

    if not candidates:
        return user_input

    file_contents = []
    remaining_text = user_input

    for path in candidates:
        if not os.path.exists(path):
            continue

        # 从用户输入中移除路径（包括引号）
        remaining_text = remaining_text.replace(f'"{path}"', '').replace(f"'{path}'", '').replace(path, '')
        filename = os.path.basename(path)
        ext = os.path.splitext(path)[1].lower()

        if os.path.isdir(path):
            # 文件夹：列出结构
            try:
                entries = os.listdir(path)[:50]
                tree = "\n".join(f"  {e}" for e in entries)
                file_contents.append(f"[文件夹: {path}]\n{tree}")
                app.write_system(f"已读取文件夹: {filename}/")
            except Exception as e:
                file_contents.append(f"[文件夹读取失败: {path}] {e}")

        elif ext in ('.png', '.jpg', '.jpeg', '.bmp', '.webp'):
            # 图片：OCR 识别
            try:
                from rapidocr_onnxruntime import RapidOCR
                ocr = RapidOCR()
                result, _ = ocr(path)
                if result:
                    text = "\n".join(line[1] for line in result)
                    file_contents.append(f"[图片 OCR: {filename}]\n{text}")
                else:
                    file_contents.append(f"[图片: {filename}]（OCR 未识别到文字）")
                app.write_system(f"已 OCR 识别: {filename}")
            except Exception as e:
                file_contents.append(f"[图片读取失败: {filename}] {e}")

        elif ext == '.pdf':
            # PDF
            try:
                import pdfplumber
                text_parts = []
                with pdfplumber.open(path) as pdf:
                    for page in pdf.pages[:20]:
                        t = page.extract_text()
                        if t:
                            text_parts.append(t)
                file_contents.append(f"[PDF: {filename}]\n" + "\n".join(text_parts))
                app.write_system(f"已读取: {filename}")
            except Exception as e:
                file_contents.append(f"[PDF 读取失败: {filename}] {e}")

        elif ext == '.docx':
            # Word 文档
            try:
                from docx import Document
                doc = Document(path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text)
                file_contents.append(f"[Word: {filename}]\n{text}")
                app.write_system(f"已读取: {filename}")
            except Exception as e:
                file_contents.append(f"[Word 读取失败: {filename}] {e}")

        elif ext == '.pptx':
            # PPT
            try:
                from pptx import Presentation
                prs = Presentation(path)
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                    if texts:
                        slides_text.append(f"第{i}页: " + " | ".join(texts))
                file_contents.append(f"[PPT: {filename}]\n" + "\n".join(slides_text))
                app.write_system(f"已读取: {filename}")
            except Exception as e:
                file_contents.append(f"[PPT 读取失败: {filename}] {e}")

        else:
            # 其他文本文件
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read(50000)
                file_contents.append(f"[文件: {filename}]\n{content}")
                app.write_system(f"已读取: {filename}")
            except Exception as e:
                file_contents.append(f"[文件读取失败: {filename}] {e}")

    if not file_contents:
        return user_input

    # 组合：用户附加文字 + 文件内容
    remaining_text = remaining_text.strip()
    parts = []
    if remaining_text:
        parts.append(remaining_text)
    parts.extend(file_contents)
    return "\n\n".join(parts)


if __name__ == "__main__":
    main()
